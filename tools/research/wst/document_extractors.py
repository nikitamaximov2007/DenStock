from __future__ import annotations

import json
from html.parser import HTMLParser
from pathlib import Path
from typing import Any


def document_source_ref(post_id: int, file_name: str, **location: Any) -> str:
    query = "&".join(f"{key}={value}" for key, value in location.items() if value is not None)
    suffix = f"?{query}" if query else ""
    return f"wst://post/{post_id}/file/{file_name}{suffix}"


def extract_document(path: Path, post_id: int) -> dict[str, Any]:
    """Extract only source text and structural positions; unsupported files remain explicit."""
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _extract_pdf(path, post_id)
    if suffix == ".pptx":
        return _extract_pptx(path, post_id)
    if suffix == ".docx":
        return _extract_docx(path, post_id)
    if suffix == ".xlsx":
        return _extract_xlsx(path, post_id)
    if suffix in {".html", ".htm"}:
        return _extract_html(path, post_id)
    if suffix in {".txt", ".md"}:
        return _result(
            path,
            post_id,
            [{"text": path.read_text(encoding="utf-8", errors="replace")}],
            "plain_text",
        )
    return _result(
        path, post_id, [], "unsupported", [f"Unsupported document format: {suffix or 'none'}"]
    )


def write_document_extraction(result: dict[str, Any], output_dir: Path) -> tuple[Path, Path]:
    output_dir.mkdir(parents=True, exist_ok=True)
    stem = f"post-{result['post_id']}-{Path(result['file_name']).stem}"
    json_path = output_dir / f"{stem}.json"
    markdown_path = output_dir / f"{stem}.md"
    json_path.write_text(json.dumps(result, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    lines = [f"# {result['file_name']}", "", f"Telegram post: {result['post_id']}", ""]
    for block in result["blocks"]:
        position = block.get("page") or block.get("slide") or block.get("sheet") or "source"
        lines.extend([f"## {position}", "", block.get("text", ""), "", block["source_ref"], ""])
    if result["errors"]:
        lines.extend(["## Processing notes", "", *(f"- {item}" for item in result["errors"]), ""])
    markdown_path.write_text("\n".join(lines), encoding="utf-8")
    return json_path, markdown_path


def _result(
    path: Path,
    post_id: int,
    raw_blocks: list[dict[str, Any]],
    method: str,
    errors: list[str] | None = None,
) -> dict[str, Any]:
    blocks = []
    for index, block in enumerate(raw_blocks, start=1):
        location = {key: block[key] for key in ("page", "slide", "sheet") if key in block}
        blocks.append(
            {
                "block_index": index,
                "text": str(block.get("text", "")).strip(),
                **location,
                "extraction_method": block.get("extraction_method", method),
                "source_ref": document_source_ref(post_id, path.name, **location),
            }
        )
    return {
        "post_id": post_id,
        "file_name": path.name,
        "source_path": str(path),
        "extraction_method": method,
        "blocks": blocks,
        "errors": errors or [],
    }


def _extract_pdf(path: Path, post_id: int) -> dict[str, Any]:
    try:
        import fitz
    except ImportError:
        return _result(path, post_id, [], "unavailable", ["PyMuPDF is not installed."])
    blocks = []
    with fitz.open(path) as document:
        for page_number, page in enumerate(document, start=1):
            text = page.get_text("text").strip()
            blocks.append(
                {
                    "page": page_number,
                    "text": text,
                    "extraction_method": "text_layer" if text else "ocr_required",
                }
            )
    method = "text_layer" if all(block["text"] for block in blocks) else "mixed"
    return _result(path, post_id, blocks, method)


def _extract_pptx(path: Path, post_id: int) -> dict[str, Any]:
    try:
        from pptx import Presentation
    except ImportError:
        return _result(path, post_id, [], "unavailable", ["python-pptx is not installed."])
    presentation = Presentation(path)
    blocks = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text:
                parts.append(shape.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    parts.append(" | ".join(cell.text for cell in row.cells))
        notes = getattr(getattr(slide, "notes_slide", None), "notes_text_frame", None)
        if notes and notes.text:
            parts.append(f"Notes: {notes.text}")
        blocks.append({"slide": slide_number, "text": "\n".join(parts)})
    return _result(path, post_id, blocks, "pptx")


def _extract_docx(path: Path, post_id: int) -> dict[str, Any]:
    try:
        from docx import Document
    except ImportError:
        return _result(path, post_id, [], "unavailable", ["python-docx is not installed."])
    document = Document(path)
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text]
    for table in document.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text for cell in row.cells))
    return _result(path, post_id, [{"text": "\n".join(parts)}], "docx")


def _extract_xlsx(path: Path, post_id: int) -> dict[str, Any]:
    try:
        from openpyxl import load_workbook
    except ImportError:
        return _result(path, post_id, [], "unavailable", ["openpyxl is not installed."])
    workbook = load_workbook(path, data_only=False, read_only=True)
    blocks = []
    for sheet in workbook.worksheets:
        rows = []
        for row in sheet.iter_rows():
            values = [
                f"{cell.coordinate}={cell.value}" for cell in row if cell.value not in (None, "")
            ]
            if values:
                rows.append(" | ".join(values))
        blocks.append({"sheet": sheet.title, "text": "\n".join(rows)})
    return _result(path, post_id, blocks, "xlsx")


class _VisibleHTML(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []
        self.skip_depth = 0

    def handle_starttag(self, tag: str, _attrs: list[tuple[str, str | None]]) -> None:
        if tag in {"script", "style", "noscript"}:
            self.skip_depth += 1

    def handle_endtag(self, tag: str) -> None:
        if tag in {"script", "style", "noscript"} and self.skip_depth:
            self.skip_depth -= 1

    def handle_data(self, data: str) -> None:
        if not self.skip_depth and data.strip():
            self.parts.append(data.strip())


def _extract_html(path: Path, post_id: int) -> dict[str, Any]:
    parser = _VisibleHTML()
    parser.feed(path.read_text(encoding="utf-8", errors="replace"))
    return _result(path, post_id, [{"text": "\n".join(parser.parts)}], "html_visible_text")
